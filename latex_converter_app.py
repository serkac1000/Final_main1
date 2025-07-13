
#!/usr/bin/env python3
import tkinter as tk
from tkinter import filedialog, messagebox, ttk
import os
import subprocess
import sys
import shutil
from pathlib import Path
import re

class LatexConverterApp:
    def __init__(self, root):
        self.root = root
        self.root.title("LaTeX to Presentation Converter")
        self.root.geometry("800x700")
        self.root.configure(bg='#f0f0f0')
        
        # Variables
        self.latex_file = tk.StringVar()
        self.output_format = tk.StringVar(value="pptx")
        self.language = tk.StringVar(value="english")
        self.media_files = []
        
        self.setup_ui()
        
    def setup_ui(self):
        # Main frame
        main_frame = tk.Frame(self.root, bg='#f0f0f0', padx=20, pady=20)
        main_frame.pack(fill='both', expand=True)
        
        # Title
        title_label = tk.Label(main_frame, 
                              text="📄 LaTeX to Presentation Converter", 
                              font=('Arial', 18, 'bold'),
                              bg='#f0f0f0',
                              fg='#0066CC')
        title_label.pack(pady=(0, 20))
        
        # File selection section
        file_frame = tk.LabelFrame(main_frame, text="📁 Input LaTeX File", 
                                  font=('Arial', 12, 'bold'),
                                  bg='#f0f0f0', fg='#333')
        file_frame.pack(fill='x', pady=(0, 15))
        
        file_select_frame = tk.Frame(file_frame, bg='#f0f0f0')
        file_select_frame.pack(fill='x', padx=10, pady=10)
        
        tk.Entry(file_select_frame, textvariable=self.latex_file, 
                font=('Arial', 10), width=60).pack(side='left', padx=(0, 10))
        
        tk.Button(file_select_frame, text="Browse", 
                 command=self.browse_latex_file,
                 bg='#0066CC', fg='white', font=('Arial', 10)).pack(side='right')
        
        # Media files section
        media_frame = tk.LabelFrame(main_frame, text="🎬 Media Files (Images/Videos)", 
                                   font=('Arial', 12, 'bold'),
                                   bg='#f0f0f0', fg='#333')
        media_frame.pack(fill='both', expand=True, pady=(0, 15))
        
        media_buttons_frame = tk.Frame(media_frame, bg='#f0f0f0')
        media_buttons_frame.pack(fill='x', padx=10, pady=5)
        
        tk.Button(media_buttons_frame, text="📸 Add Images", 
                 command=self.add_images,
                 bg='#28a745', fg='white', font=('Arial', 10)).pack(side='left', padx=(0, 10))
        
        tk.Button(media_buttons_frame, text="🎥 Add Videos", 
                 command=self.add_videos,
                 bg='#dc3545', fg='white', font=('Arial', 10)).pack(side='left', padx=(0, 10))
        
        tk.Button(media_buttons_frame, text="🗑️ Clear All", 
                 command=self.clear_media,
                 bg='#6c757d', fg='white', font=('Arial', 10)).pack(side='right')
        
        # Media list
        self.media_listbox = tk.Listbox(media_frame, height=8, font=('Arial', 9))
        self.media_listbox.pack(fill='both', expand=True, padx=10, pady=(5, 10))
        
        scrollbar = tk.Scrollbar(media_frame)
        scrollbar.pack(side='right', fill='y')
        self.media_listbox.config(yscrollcommand=scrollbar.set)
        scrollbar.config(command=self.media_listbox.yview)
        
        # Language selection
        lang_frame = tk.LabelFrame(main_frame, text="🌍 Language", 
                                  font=('Arial', 12, 'bold'),
                                  bg='#f0f0f0', fg='#333')
        lang_frame.pack(fill='x', pady=(0, 15))
        
        lang_options_frame = tk.Frame(lang_frame, bg='#f0f0f0')
        lang_options_frame.pack(padx=10, pady=10)
        
        tk.Radiobutton(lang_options_frame, text="🇺🇸 English", 
                      variable=self.language, value="english",
                      font=('Arial', 11), bg='#f0f0f0').pack(side='left', padx=(0, 20))
        
        tk.Radiobutton(lang_options_frame, text="🇷🇺 Russian", 
                      variable=self.language, value="russian",
                      font=('Arial', 11), bg='#f0f0f0').pack(side='left')
        
        # Output format selection
        format_frame = tk.LabelFrame(main_frame, text="📊 Output Format", 
                                    font=('Arial', 12, 'bold'),
                                    bg='#f0f0f0', fg='#333')
        format_frame.pack(fill='x', pady=(0, 15))
        
        format_options_frame = tk.Frame(format_frame, bg='#f0f0f0')
        format_options_frame.pack(padx=10, pady=10)
        
        tk.Radiobutton(format_options_frame, text="📊 PowerPoint (.pptx)", 
                      variable=self.output_format, value="pptx",
                      font=('Arial', 11), bg='#f0f0f0').pack(side='left', padx=(0, 20))
        
        tk.Radiobutton(format_options_frame, text="📄 PDF Document (.pdf)", 
                      variable=self.output_format, value="pdf",
                      font=('Arial', 11), bg='#f0f0f0').pack(side='left')
        
        # Progress and status
        self.status_label = tk.Label(main_frame, text="Ready to convert...", 
                                    font=('Arial', 10), fg='green', bg='#f0f0f0')
        self.status_label.pack(pady=(10, 5))
        
        self.progress = ttk.Progressbar(main_frame, mode='indeterminate')
        self.progress.pack(fill='x', pady=(0, 15))
        
        # Convert button
        convert_button = tk.Button(main_frame, text="🚀 Convert to Presentation", 
                                  command=self.convert_presentation,
                                  font=('Arial', 14, 'bold'),
                                  bg='#0066CC', fg='white',
                                  padx=20, pady=10)
        convert_button.pack(pady=10)
        
    def browse_latex_file(self):
        file_path = filedialog.askopenfilename(
            title="Select LaTeX File",
            filetypes=[("LaTeX files", "*.tex"), ("All files", "*.*")]
        )
        if file_path:
            self.latex_file.set(file_path)
            
    def add_images(self):
        files = filedialog.askopenfilenames(
            title="Select Image Files",
            filetypes=[
                ("Image files", "*.png *.jpg *.jpeg *.gif *.bmp *.tiff"),
                ("All files", "*.*")
            ]
        )
        for file in files:
            if file not in self.media_files:
                self.media_files.append(file)
                self.media_listbox.insert(tk.END, f"📸 {os.path.basename(file)}")
                
    def add_videos(self):
        files = filedialog.askopenfilenames(
            title="Select Video Files",
            filetypes=[
                ("Video files", "*.mp4 *.avi *.mov *.wmv *.flv *.mkv"),
                ("All files", "*.*")
            ]
        )
        for file in files:
            if file not in self.media_files:
                self.media_files.append(file)
                self.media_listbox.insert(tk.END, f"🎥 {os.path.basename(file)}")
                
    def clear_media(self):
        self.media_files.clear()
        self.media_listbox.delete(0, tk.END)
        
    def parse_latex_content(self, latex_file):
        """Parse LaTeX file and extract content for presentation"""
        try:
            with open(latex_file, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # Extract title
            title_match = re.search(r'\\title\{([^}]+)\}', content)
            title = title_match.group(1) if title_match else "Presentation"
            
            # Extract sections
            sections = re.findall(r'\\section\{([^}]+)\}', content)
            
            # Extract subsections with content
            subsections = []
            subsection_pattern = r'\\subsection\{([^}]+)\}(.*?)(?=\\subsection|\\section|\\end\{document\})'
            matches = re.findall(subsection_pattern, content, re.DOTALL)
            
            for match in matches:
                subsection_title = match[0]
                subsection_content = match[1].strip()
                
                # Extract bullet points
                bullets = re.findall(r'\\item\s+([^\n\\]+)', subsection_content)
                
                subsections.append({
                    'title': subsection_title,
                    'content': subsection_content,
                    'bullets': bullets
                })
            
            return {
                'title': title,
                'sections': sections,
                'subsections': subsections
            }
            
        except Exception as e:
            raise Exception(f"Error parsing LaTeX file: {str(e)}")
            
    def create_presentation_from_latex(self, latex_data):
        """Create presentation file from parsed LaTeX data"""
        output_dir = "output_presentations"
        os.makedirs(output_dir, exist_ok=True)
        
        # Copy media files to output directory
        media_dir = os.path.join(output_dir, "media")
        os.makedirs(media_dir, exist_ok=True)
        
        copied_media = []
        for media_file in self.media_files:
            dest_path = os.path.join(media_dir, os.path.basename(media_file))
            shutil.copy2(media_file, dest_path)
            copied_media.append(dest_path)
        
        if self.output_format.get() == "pptx":
            return self.create_pptx_from_latex(latex_data, output_dir, copied_media)
        else:
            return self.create_pdf_from_latex(latex_data, output_dir, copied_media)
            
    def create_pptx_from_latex(self, latex_data, output_dir, media_files):
        """Create PowerPoint presentation from LaTeX data"""
        try:
            from pptx import Presentation
            from pptx.util import Inches
            from pptx.dml.color import RGBColor
            
            prs = Presentation()
            
            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = latex_data['title']
            subtitle.text = f"Generated from LaTeX\n{len(media_files)} media files included"
            
            # Content slides
            bullet_slide_layout = prs.slide_layouts[1]
            
            for subsection in latex_data['subsections']:
                slide = prs.slides.add_slide(bullet_slide_layout)
                title = slide.shapes.title
                content = slide.placeholders[1]
                
                title.text = subsection['title']
                tf = content.text_frame
                
                if subsection['bullets']:
                    tf.text = subsection['bullets'][0] if subsection['bullets'] else "Content"
                    for bullet in subsection['bullets'][1:]:
                        p = tf.add_paragraph()
                        p.text = bullet
                        p.level = 0
                else:
                    # Use raw content if no bullets found
                    clean_content = re.sub(r'\\[a-zA-Z]+\{[^}]*\}', '', subsection['content'])
                    clean_content = re.sub(r'[{}\\]', '', clean_content).strip()
                    tf.text = clean_content[:200] + "..." if len(clean_content) > 200 else clean_content
                
                # Add media if available
                if media_files:
                    media_file = media_files[len(prs.slides) % len(media_files) - 1]
                    if media_file.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp')):
                        try:
                            slide.shapes.add_picture(media_file, 
                                                   Inches(6), Inches(2), 
                                                   Inches(3), Inches(2))
                        except:
                            pass  # Skip if image can't be added
            
            # Save presentation
            lang_suffix = "_russian" if self.language.get() == "russian" else "_english"
            filename = os.path.join(output_dir, f"latex_presentation{lang_suffix}.pptx")
            prs.save(filename)
            
            return filename
            
        except ImportError:
            raise Exception("python-pptx library not found. Installing...")
        except Exception as e:
            raise Exception(f"Error creating PowerPoint: {str(e)}")
            
    def create_pdf_from_latex(self, latex_data, output_dir, media_files):
        """Create PDF presentation from LaTeX data"""
        try:
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter, A4
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, PageBreak
            from reportlab.lib.units import inch
            from reportlab.lib.colors import HexColor
            from reportlab.lib.enums import TA_CENTER, TA_LEFT
            
            lang_suffix = "_russian" if self.language.get() == "russian" else "_english"
            filename = os.path.join(output_dir, f"latex_presentation{lang_suffix}.pdf")
            
            doc = SimpleDocTemplate(filename, pagesize=A4)
            styles = getSampleStyleSheet()
            
            # Custom styles
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Title'],
                fontSize=24,
                spaceAfter=30,
                alignment=TA_CENTER,
                textColor=HexColor('#0066CC')
            )
            
            heading_style = ParagraphStyle(
                'CustomHeading',
                parent=styles['Heading1'],
                fontSize=18,
                spaceAfter=20,
                textColor=HexColor('#0066CC')
            )
            
            story = []
            
            # Title page
            story.append(Paragraph(latex_data['title'], title_style))
            story.append(Spacer(1, 0.5*inch))
            story.append(Paragraph("Generated from LaTeX", heading_style))
            story.append(Spacer(1, 0.3*inch))
            story.append(Paragraph(f"{len(media_files)} media files included", styles['Normal']))
            story.append(PageBreak())
            
            # Content pages
            for i, subsection in enumerate(latex_data['subsections']):
                story.append(Paragraph(subsection['title'], heading_style))
                
                if subsection['bullets']:
                    for bullet in subsection['bullets']:
                        story.append(Paragraph(f"• {bullet}", styles['Normal']))
                        story.append(Spacer(1, 6))
                else:
                    clean_content = re.sub(r'\\[a-zA-Z]+\{[^}]*\}', '', subsection['content'])
                    clean_content = re.sub(r'[{}\\]', '', clean_content).strip()
                    if clean_content:
                        story.append(Paragraph(clean_content, styles['Normal']))
                
                # Add image if available
                if media_files and i < len(media_files):
                    media_file = media_files[i]
                    if media_file.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp')):
                        try:
                            story.append(Spacer(1, 0.2*inch))
                            img = Image(media_file, width=4*inch, height=3*inch)
                            story.append(img)
                        except:
                            pass
                
                story.append(PageBreak())
            
            doc.build(story)
            return filename
            
        except ImportError:
            raise Exception("reportlab library not found. Installing...")
        except Exception as e:
            raise Exception(f"Error creating PDF: {str(e)}")
            
    def convert_presentation(self):
        if not self.latex_file.get():
            messagebox.showerror("Error", "Please select a LaTeX file first!")
            return
            
        if not os.path.exists(self.latex_file.get()):
            messagebox.showerror("Error", "Selected LaTeX file does not exist!")
            return
            
        try:
            self.status_label.config(text="Converting...", fg='blue')
            self.progress.start()
            self.root.update()
            
            # Parse LaTeX content
            latex_data = self.parse_latex_content(self.latex_file.get())
            
            # Create presentation
            output_file = self.create_presentation_from_latex(latex_data)
            
            self.progress.stop()
            self.status_label.config(text="Conversion completed successfully!", fg='green')
            
            # Try to open the file
            try:
                if os.name == 'nt':  # Windows
                    os.startfile(output_file)
                elif os.name == 'posix':  # macOS and Linux
                    subprocess.run(['open' if sys.platform == 'darwin' else 'xdg-open', output_file])
                    
                messagebox.showinfo("Success", 
                                  f"Presentation created successfully!\n\nFile: {output_file}\n\nMedia files: {len(self.media_files)}")
            except:
                messagebox.showinfo("Success", 
                                  f"Presentation created successfully!\n\nFile: {output_file}\n\nPlease open it manually.")
                
        except Exception as e:
            self.progress.stop()
            self.status_label.config(text="Error occurred during conversion", fg='red')
            messagebox.showerror("Error", f"Conversion failed:\n{str(e)}")

def main():
    root = tk.Tk()
    app = LatexConverterApp(root)
    root.mainloop()

if __name__ == "__main__":
    main()
