
#!/usr/bin/env python3
import os
import json
import tempfile
import subprocess
from datetime import datetime
from flask import Flask, request, jsonify, send_file, render_template_string
from werkzeug.utils import secure_filename
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.colors import HexColor
from reportlab.lib.enums import TA_CENTER, TA_LEFT

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'temp_files'
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB max

# Create temp directory
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

class LaTeXToPresentationConverter:
    def __init__(self):
        self.colors = {
            'primary': RGBColor(0, 102, 204),
            'secondary': RGBColor(255, 153, 51),
            'success': RGBColor(40, 167, 69),
            'danger': RGBColor(220, 53, 69)
        }

    def parse_latex(self, latex_code):
        """Parse LaTeX code and extract presentation structure"""
        slides = []
        current_slide = None
        
        lines = latex_code.split('\n')
        in_frame = False
        in_itemize = False
        frame_title = ""
        frame_content = []
        
        for line in lines:
            line = line.strip()
            
            # Extract document title
            if line.startswith('\\title{'):
                title = line.replace('\\title{', '').replace('}', '')
                slides.append({'type': 'title', 'title': title, 'content': []})
            
            # Extract author
            elif line.startswith('\\author{'):
                author = line.replace('\\author{', '').replace('}', '')
                if slides and slides[-1]['type'] == 'title':
                    slides[-1]['author'] = author
            
            # Frame start
            elif line.startswith('\\begin{frame}'):
                in_frame = True
                frame_content = []
                # Extract frame title
                if '{' in line:
                    frame_title = line.split('{', 2)[-1].replace('}', '')
                else:
                    frame_title = "Slide"
            
            # Frame end
            elif line.startswith('\\end{frame}'):
                if in_frame:
                    slides.append({
                        'type': 'content',
                        'title': frame_title,
                        'content': frame_content.copy()
                    })
                in_frame = False
                in_itemize = False
            
            # Itemize start
            elif line.startswith('\\begin{itemize}') and in_frame:
                in_itemize = True
            
            # Itemize end
            elif line.startswith('\\end{itemize}') and in_frame:
                in_itemize = False
            
            # List items
            elif line.startswith('\\item') and in_frame and in_itemize:
                item_text = line.replace('\\item', '').strip()
                frame_content.append({'type': 'bullet', 'text': item_text})
            
            # Regular text in frame
            elif in_frame and line and not line.startswith('\\'):
                frame_content.append({'type': 'text', 'text': line})
        
        return slides

    def create_pptx_presentation(self, slides, language='english'):
        """Create PowerPoint presentation from parsed slides"""
        prs = Presentation()
        
        # Title slide
        if slides and slides[0]['type'] == 'title':
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_slide.shapes.title.text = slides[0]['title']
            if 'author' in slides[0]:
                title_slide.placeholders[1].text = slides[0]['author']
            slides = slides[1:]  # Remove title slide from content slides
        
        # Content slides
        for slide_data in slides:
            if slide_data['type'] == 'content':
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = slide_data['title']
                
                if slide_data['content']:
                    content_placeholder = slide.placeholders[1]
                    tf = content_placeholder.text_frame
                    tf.clear()
                    
                    for i, item in enumerate(slide_data['content']):
                        if i == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()
                        
                        p.text = item['text']
                        if item['type'] == 'bullet':
                            p.level = 0
        
        # Save presentation
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        lang_suffix = "_Russian" if language == 'russian' else "_English"
        filename = f"LaTeX_Presentation{lang_suffix}_{timestamp}.pptx"
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        prs.save(filepath)
        
        return filename

    def create_pdf_presentation(self, slides, language='english'):
        """Create PDF presentation from parsed slides"""
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        lang_suffix = "_Russian" if language == 'russian' else "_English"
        filename = f"LaTeX_Presentation{lang_suffix}_{timestamp}.pdf"
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        
        doc = SimpleDocTemplate(filepath, pagesize=A4)
        styles = getSampleStyleSheet()
        
        # Custom styles
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Title'],
            fontSize=24,
            spaceAfter=30,
            alignment=TA_CENTER,
            textColor=HexColor('#0066CC')
        )
        
        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading1'],
            fontSize=18,
            spaceAfter=20,
            textColor=HexColor('#0066CC')
        )
        
        bullet_style = ParagraphStyle(
            'BulletText',
            parent=styles['Normal'],
            fontSize=12,
            spaceAfter=8,
            leftIndent=20
        )
        
        story = []
        
        # Title page
        if slides and slides[0]['type'] == 'title':
            story.append(Paragraph(slides[0]['title'], title_style))
            if 'author' in slides[0]:
                story.append(Spacer(1, 0.5*72))
                story.append(Paragraph(slides[0]['author'], styles['Normal']))
            story.append(PageBreak())
            slides = slides[1:]
        
        # Content slides
        for slide_data in slides:
            if slide_data['type'] == 'content':
                story.append(Paragraph(slide_data['title'], heading_style))
                
                for item in slide_data['content']:
                    if item['type'] == 'bullet':
                        story.append(Paragraph(f"• {item['text']}", bullet_style))
                    else:
                        story.append(Paragraph(item['text'], styles['Normal']))
                
                story.append(PageBreak())
        
        doc.build(story)
        return filename

    def translate_content(self, slides, target_language):
        """Basic translation mapping for common terms"""
        if target_language == 'russian':
            translations = {
                'Introduction': 'Введение',
                'Overview': 'Обзор',
                'Conclusion': 'Заключение',
                'Summary': 'Резюме',
                'Features': 'Функции',
                'Benefits': 'Преимущества',
                'Getting Started': 'Начало работы',
                'Technical': 'Технический',
                'Implementation': 'Реализация'
            }
            
            for slide in slides:
                if 'title' in slide:
                    for eng, rus in translations.items():
                        slide['title'] = slide['title'].replace(eng, rus)
        
        return slides

converter = LaTeXToPresentationConverter()

@app.route('/')
def index():
    with open('latex_converter.html', 'r', encoding='utf-8') as f:
        return f.read()

@app.route('/convert', methods=['POST'])
def convert_latex():
    try:
        data = request.get_json()
        latex_code = data.get('latex', '')
        language = data.get('language', 'english')
        format_type = data.get('format', 'pptx')
        
        if not latex_code:
            return jsonify({'success': False, 'error': 'No LaTeX code provided'})
        
        # Parse LaTeX
        slides = converter.parse_latex(latex_code)
        
        if not slides:
            return jsonify({'success': False, 'error': 'No slides found in LaTeX code'})
        
        # Translate if needed
        if language == 'russian':
            slides = converter.translate_content(slides, 'russian')
        
        # Convert to requested format
        if format_type == 'pptx':
            filename = converter.create_pptx_presentation(slides, language)
        else:
            filename = converter.create_pdf_presentation(slides, language)
        
        return jsonify({
            'success': True,
            'filename': filename,
            'slides_count': len(slides)
        })
        
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)})

@app.route('/download/<filename>')
def download_file(filename):
    try:
        filename = secure_filename(filename)
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        
        if os.path.exists(filepath):
            return send_file(filepath, as_attachment=True)
        else:
            return jsonify({'error': 'File not found'}), 404
            
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/cleanup')
def cleanup():
    """Clean up temporary files older than 1 hour"""
    try:
        import time
        current_time = time.time()
        count = 0
        
        for filename in os.listdir(app.config['UPLOAD_FOLDER']):
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            if os.path.isfile(filepath):
                file_time = os.path.getctime(filepath)
                if current_time - file_time > 3600:  # 1 hour
                    os.remove(filepath)
                    count += 1
        
        return jsonify({'cleaned': count, 'message': f'Cleaned {count} old files'})
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

if __name__ == '__main__':
    print("🚀 LaTeX to Presentation Converter Server Starting...")
    print("📊 Supports: PPTX and PDF output")
    print("🌍 Languages: English and Russian")
    print("🔗 Access at: http://0.0.0.0:5000")
    app.run(host='0.0.0.0', port=5000, debug=True)
